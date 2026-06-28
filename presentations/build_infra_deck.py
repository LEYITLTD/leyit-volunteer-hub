#!/usr/bin/env python3
"""Build the technical infrastructure deck (Architecture / Data model / Auth)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ── Brand ─────────────────────────────────────────────────────────────────
GOLD  = RGBColor(0xA8,0x85,0x4A); GOLDL = RGBColor(0xF3,0xEF,0xE6)
INK   = RGBColor(0x1A,0x17,0x14); PAPER = RGBColor(0xF7,0xF4,0xEE)
CARD  = RGBColor(0xFF,0xFF,0xFF); MUTED = RGBColor(0x78,0x71,0x6C)
LINE  = RGBColor(0xE4,0xDD,0xD0); WHITE = RGBColor(0xFF,0xFF,0xFF)
BLUE  = RGBColor(0x1D,0x4E,0xD8); BLUEL = RGBColor(0xE7,0xEE,0xF6)
GREEN = RGBColor(0x15,0x80,0x3D); GREENL= RGBColor(0xDC,0xFC,0xE7)
RED   = RGBColor(0xDC,0x26,0x26); REDL  = RGBColor(0xFE,0xE2,0xE2)
PURP  = RGBColor(0x5B,0x21,0xB6); PURPL = RGBColor(0xED,0xE9,0xFE)
WARN  = RGBColor(0x92,0x40,0x0E); WARNL = RGBColor(0xFE,0xF9,0xC3)
FONT  = "Calibri"

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
W, H = 13.333, 7.5

def slide(bg=PAPER):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = bg
    return s

def _noline(sh): sh.line.fill.background()

def box(s, x, y, w, h, fill, line=None, rounded=True, lw=1.0):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: _noline(shp)
    else: shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(1)
    for i, ln in enumerate(runs):
        txt, size, bold, color = ln[0], ln[1], ln[2], ln[3]
        sa = ln[4] if len(ln) > 4 else 3
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sa); p.space_before = Pt(0)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = FONT
    return tb

def boxlabel(s, x, y, w, h, title, sub, fill, tc, sc=None, line=None, tsize=12, ssize=9):
    box(s, x, y, w, h, fill, line=line, lw=1.25)
    runs = [(title, tsize, True, tc, 2)]
    if sub: runs.append((sub, ssize, False, sc or tc, 0))
    text(s, x+0.08, y, w-0.16, h, runs, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def arrow(s, x1, y1, x2, y2, color=GOLD, w=1.75, dash=False):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(w)
    try:
        ln = c.line._get_or_add_ln()
        tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
        ln.append(tail)
        if dash:
            d = ln.makeelement(qn('a:prstDash'), {'val': 'dash'}); ln.insert(0, d)
    except Exception:
        pass
    return c

def header(s, tag, title, sub=None):
    box(s, 0, 0, W, 0.16, GOLD, rounded=False)
    text(s, 0.6, 0.42, 12, 0.35, [(tag.upper(), 12, True, GOLD, 0)])
    text(s, 0.6, 0.74, 12.1, 0.7, [(title, 27, True, INK, 0)])
    if sub: text(s, 0.62, 1.32, 12.1, 0.4, [(sub, 13, False, MUTED, 0)])
    return 1.85 if sub else 1.7

def footer(s, n):
    text(s, 0.6, 7.06, 8, 0.3, [("Volunteer Hub · Technical Infrastructure", 9, False, MUTED, 0)])
    text(s, 12.0, 7.06, 0.9, 0.3, [(str(n), 9, False, MUTED, 0)], align=PP_ALIGN.RIGHT)

PAGE = [0]
def page(s):
    PAGE[0] += 1; footer(s, PAGE[0])

# ── Field table ───────────────────────────────────────────────────────────
def field_table(s, x, y, w, title, fields, accent=GOLD, h_each=0.305, title_h=0.34):
    text(s, x, y, w, title_h, [(title, 13, True, INK, 0)])
    ty = y + title_h
    rows = len(fields) + 1
    tbl_h = h_each * rows
    gfx = s.shapes.add_table(rows, 2, Inches(x), Inches(ty), Inches(w), Inches(tbl_h))
    t = gfx.table
    t.first_row = False; t.horz_banding = False
    t.columns[0].width = Inches(w*0.42); t.columns[1].width = Inches(w*0.58)
    # header row
    for c, lab in enumerate(["Column", "Type / notes"]):
        cell = t.cell(0, c); cell.fill.solid(); cell.fill.fore_color.rgb = accent
        cell.margin_left = Pt(6); cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
        p = cell.text_frame.paragraphs[0]; r = p.add_run(); r.text = lab
        r.font.size = Pt(9.5); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
    for i, (col, note) in enumerate(fields):
        for c, val, bold, color, size in [(0, col, True, INK, 9.7), (1, note, False, MUTED, 9.2)]:
            cell = t.cell(i+1, c)
            cell.fill.solid(); cell.fill.fore_color.rgb = CARD if i % 2 == 0 else GOLDL
            cell.margin_left = Pt(6); cell.margin_top = Pt(1); cell.margin_bottom = Pt(1)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]; rr = p.add_run(); rr.text = val
            rr.font.size = Pt(size); rr.font.bold = bold; rr.font.color.rgb = color; rr.font.name = FONT
    return ty + tbl_h

def bullets(s, x, y, w, items, size=13, gap=7, color=INK):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        lead, body = (it if isinstance(it, tuple) else (None, it))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.space_before = Pt(0)
        d = p.add_run(); d.text = "•  "; d.font.size = Pt(size); d.font.color.rgb = GOLD; d.font.bold = True; d.font.name = FONT
        if lead:
            r = p.add_run(); r.text = lead + "  "; r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = color; r.font.name = FONT
        r2 = p.add_run(); r2.text = body; r2.font.size = Pt(size); r2.font.color.rgb = MUTED if lead else color; r2.font.name = FONT
    return tb

def callout(s, x, y, w, h, title, body, fill=WARNL, tc=WARN):
    box(s, x, y, w, h, fill, rounded=True)
    text(s, x+0.2, y+0.12, w-0.4, h-0.2, [(title, 11.5, True, tc, 3), (body, 11, False, tc, 0)], anchor=MSO_ANCHOR.TOP)

# ════════════════════════════════════════════════════════════════════════
# TITLE
# ════════════════════════════════════════════════════════════════════════
s = slide(INK)
box(s, 0, 3.05, W, 0.06, GOLD, rounded=False)
text(s, 0.9, 1.7, 11.5, 1.0, [("Volunteer Hub", 16, True, GOLD, 0)])
text(s, 0.9, 2.15, 11.5, 1.2, [("Technical Infrastructure", 44, True, WHITE, 0)])
text(s, 0.92, 3.35, 11.5, 1.0, [("Architecture  ·  Data Model  ·  Authentication & Onboarding", 17, False, RGBColor(0xC9,0xC1,0xB2), 0)])
text(s, 0.92, 6.5, 11.5, 0.5, [("Next.js 16 · React 19 · Supabase (Postgres / Auth / Storage) · Vercel", 12, False, MUTED, 0)])

def divider(num, title, sub):
    s = slide(INK)
    box(s, 0.9, 2.8, 0.9, 0.9, GOLD)
    text(s, 0.9, 2.8, 0.9, 0.9, [(num, 30, True, INK, 0)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 2.1, 2.75, 10.5, 1.0, [(title, 36, True, WHITE, 0)])
    text(s, 2.12, 3.75, 10.5, 0.6, [(sub, 15, False, RGBColor(0xC9,0xC1,0xB2), 0)])
    page(s)

# ════════════════════════════════════════════════════════════════════════
# 1 · ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════
divider("1", "Architecture", "How the system is composed and how a request flows through it")

# 1a — System context
s = slide(); header(s, "Architecture · 1 of 5", "System context")
# clients
boxlabel(s, 0.6, 2.0, 2.5, 0.95, "Volunteers", "Mobile-first web", BLUEL, BLUE, BLUE, line=BLUE)
boxlabel(s, 0.6, 3.25, 2.5, 0.95, "Admins", "Dashboard + OTP", BLUEL, BLUE, BLUE, line=BLUE)
# app
boxlabel(s, 4.2, 2.35, 3.2, 2.0, "Next.js 16 app", "App Router · server route\nhandlers · React 19 UI\n(hosted on Vercel)", INK, WHITE, RGBColor(0xC9,0xC1,0xB2))
# supabase
boxlabel(s, 8.5, 1.7, 4.2, 1.25, "Supabase — Postgres", "Row-level security · triggers", GREENL, GREEN, GREEN, line=GREEN)
boxlabel(s, 8.5, 3.05, 2.0, 1.0, "Auth", "JWT sessions", GREENL, GREEN, GREEN, line=GREEN)
boxlabel(s, 10.7, 3.05, 2.0, 1.0, "Storage", "DBS · certs", GREENL, GREEN, GREEN, line=GREEN)
# external
boxlabel(s, 8.5, 4.35, 1.3, 0.85, "Resend", "Email", GOLDL, INK, MUTED, line=LINE, tsize=10, ssize=8)
boxlabel(s, 9.95, 4.35, 1.35, 0.85, "Voodoo", "SMS", GOLDL, INK, MUTED, line=LINE, tsize=10, ssize=8)
boxlabel(s, 11.45, 4.35, 1.25, 0.85, "R2", "Object store", GOLDL, INK, MUTED, line=LINE, tsize=10, ssize=8)
arrow(s, 3.1, 2.47, 4.2, 2.7); arrow(s, 3.1, 3.72, 4.2, 3.5)
arrow(s, 7.4, 2.9, 8.5, 2.4); arrow(s, 7.4, 3.3, 8.5, 3.5); arrow(s, 7.4, 3.7, 10.7, 3.5)
arrow(s, 7.4, 4.0, 8.5, 4.6, color=MUTED)
text(s, 4.2, 4.7, 3.2, 0.5, [("Webhooks (delivery receipts) ← Resend / Voodoo", 9.5, False, MUTED, 0)])
arrow(s, 9.1, 5.2, 6.2, 4.95, color=MUTED, dash=True)
callout(s, 0.6, 5.6, 12.1, 1.05, "In one line",
        "Browsers talk only to the Next.js app; the app is the single gateway to Postgres, Auth, Storage and the email/SMS providers. Providers call back via signed webhooks to record delivery status.",
        fill=GOLDL, tc=INK)
page(s)

# 1b — Stack & conventions
s = slide(); header(s, "Architecture · 2 of 5", "Runtime & conventions")
bullets(s, 0.7, 2.0, 6.0, [
    ("Next.js 16 (App Router).", "Server route handlers under /api/* are the entire backend; React Server + Client Components for the UI."),
    ("React 19 + TypeScript + Tailwind v4.", "Strongly-typed; design tokens for the gold/espresso theme."),
    ("Vercel hosting.", "Serverless functions; env-var driven config per environment."),
    ("No custom middleware.", "Every protected route enforces auth itself via a guard — single, explicit gate per handler."),
])
bullets(s, 7.0, 2.0, 5.7, [
    ("Two Supabase clients.", "An anon/SSR client (carries the user session, RLS-enforced) and a service-role client (server-only, bypasses RLS)."),
    ("Append-only & trigger-derived data.", "e.g. points are append-only; compliance status is computed by a DB trigger."),
    ("Email/SMS as templates.", "Admin-editable templates rendered + sent via Resend / Voodoo."),
])
callout(s, 0.7, 5.75, 12.0, 0.95, "Backend = route handlers",
        "There is no separate API server. Each file under app/api/* is a serverless endpoint that authenticates, then uses the appropriate Supabase client.", fill=BLUEL, tc=BLUE)
page(s)

# 1c — Request lifecycle
s = slide(); header(s, "Architecture · 3 of 5", "Request lifecycle", "A typical authenticated API call, end to end")
steps = [
    ("Browser", "fetch('/api/…')\nwith session cookie", BLUEL, BLUE),
    ("Route handler", "app/api/*/route.ts\n(serverless on Vercel)", INK, WHITE),
    ("Auth guard", "requireAdminUser()  or\nsupabase.auth.getUser()", GOLDL, INK),
    ("Supabase client", "service-role (bypass RLS)\nor anon (RLS-enforced)", GREENL, GREEN),
    ("Postgres", "query + triggers,\nreturns rows", GREENL, GREEN),
]
x = 0.6; bw = 2.3; gap = 0.32; y = 2.5
for i,(t,sub,fill,tc) in enumerate(steps):
    boxlabel(s, x, y, bw, 1.5, t, sub, fill, tc, RGBColor(0xC9,0xC1,0xB2) if fill==INK else tc, line=None if fill in (INK,) else (tc if fill in (BLUEL,GREENL) else LINE), tsize=12.5, ssize=9.5)
    if i < len(steps)-1: arrow(s, x+bw, y+0.75, x+bw+gap, y+0.75)
    x += bw + gap
arrow(s, x-gap-0.05, y+1.55, 1.75, y+1.55, color=MUTED, dash=True)
text(s, 4.6, y+1.62, 6.0, 0.4, [("JSON response ← returned to the browser", 10.5, False, MUTED, 0)])
callout(s, 0.6, 5.5, 12.1, 1.1, "Where authorization lives",
        "The guard runs before any data access. Admin endpoints check the JWT's app_metadata.role (set server-side, not user-editable); volunteer endpoints resolve the volunteer from the session — never from a client-supplied id.", fill=GOLDL, tc=INK)
page(s)

# 1d — Trust boundary: two clients
s = slide(); header(s, "Architecture · 4 of 5", "Supabase clients & the trust boundary")
boxlabel(s, 0.7, 2.2, 5.7, 0.7, "Anon / SSR client", "carries the end-user session", BLUEL, BLUE, BLUE, line=BLUE, tsize=14, ssize=10)
bullets(s, 0.85, 3.05, 5.5, [
    ("Used by", "the browser & cookie-bound server reads."),
    ("RLS enforced", "— only rows the policy allows."),
    ("Public anon key", "ships to the browser (not a secret)."),
], size=11.5, gap=6)
boxlabel(s, 6.9, 2.2, 5.7, 0.7, "Service-role client", "server-only, full access", GREENL, GREEN, GREEN, line=GREEN, tsize=14, ssize=10)
bullets(s, 7.05, 3.05, 5.4, [
    ("Used by", "route handlers for privileged writes."),
    ("Bypasses RLS", "— so handlers must authorize first."),
    ("SUPABASE_SERVICE_ROLE_KEY", "server-only; never in client code."),
], size=11.5, gap=6)
callout(s, 0.7, 5.55, 11.9, 1.05, "Why it matters",
        "Because the service-role client bypasses RLS, the security boundary is the route handler's guard. RLS is still the safety net: most tables are deny-by-default (RLS on, no anon policy), so a leaked anon key can't read them.", fill=WARNL, tc=WARN)
page(s)

# 1e — Integrations & webhooks
s = slide(); header(s, "Architecture · 5 of 5", "External integrations")
field_table(s, 0.7, 1.95, 6.1, "Outbound", [
    ("Resend", "Transactional + broadcast email"),
    ("Voodoo SMS", "SMS broadcasts + direct SMS"),
    ("Supabase Storage", "dbs-documents (private), certificates"),
    ("Cloudflare R2", "Object storage (assets)"),
], h_each=0.34)
field_table(s, 7.0, 1.95, 5.6, "Inbound webhooks", [
    ("/api/webhooks/resend", "Delivery / open / bounce — HMAC-verified"),
    ("/api/webhooks/voodoo", "SMS delivery status — token-verified"),
], h_each=0.34)
callout(s, 7.0, 3.45, 5.6, 1.6, "Delivery tracking",
        "Each broadcast recipient row stores the provider message id. Webhooks (or an on-demand poll) flip the row to delivered / failed, powering the per-recipient ‘who received / who failed’ report.", fill=BLUEL, tc=BLUE)
callout(s, 0.7, 5.5, 6.1, 1.15, "Secrets",
        "All provider keys are server-only env vars. Only NEXT_PUBLIC_SUPABASE_URL / ANON_KEY / APP_URL are exposed to the browser.", fill=GOLDL, tc=INK)
page(s)

# ════════════════════════════════════════════════════════════════════════
# 2 · DATA MODEL
# ════════════════════════════════════════════════════════════════════════
divider("2", "Data Model", "20 tables across identity, compliance, events, engagement & messaging")

# 2a — ERD overview
s = slide(); header(s, "Data model · 1 of 8", "Entity relationships (core)")
# central
boxlabel(s, 5.5, 3.15, 2.3, 0.85, "volunteers", "the person", INK, WHITE, RGBColor(0xC9,0xC1,0xB2), tsize=12.5)
boxlabel(s, 5.5, 1.65, 2.3, 0.8, "volunteer_compliance", "1 — 1", GREENL, GREEN, GREEN, line=GREEN, tsize=11, ssize=9)
boxlabel(s, 0.7, 3.2, 1.9, 0.8, "admins", "staff", GOLDL, INK, MUTED, line=LINE, tsize=11)
boxlabel(s, 9.3, 1.55, 2.0, 0.8, "events", "", GOLDL, INK, MUTED, line=LINE, tsize=11)
boxlabel(s, 11.5, 1.55, 1.6, 0.8, "roles", "catalog", GOLDL, INK, MUTED, line=LINE, tsize=10, ssize=8)
boxlabel(s, 9.3, 2.65, 2.0, 0.8, "event_roles", "", GOLDL, INK, MUTED, line=LINE, tsize=10.5)
boxlabel(s, 9.3, 3.75, 2.2, 0.8, "event_applications", "sign-ups", BLUEL, BLUE, BLUE, line=BLUE, tsize=10.5, ssize=8)
boxlabel(s, 9.3, 4.85, 2.0, 0.8, "check_ins", "QR scans", BLUEL, BLUE, BLUE, line=BLUE, tsize=10.5, ssize=8)
boxlabel(s, 5.5, 4.8, 2.3, 0.8, "points_transactions", "append-only", PURPL, PURP, PURP, line=PURP, tsize=10.5, ssize=8)
boxlabel(s, 0.7, 4.7, 2.0, 0.8, "broadcast_recipients", "", BLUEL, BLUE, BLUE, line=BLUE, tsize=9.5)
boxlabel(s, 2.95, 4.7, 1.9, 0.8, "communications", "per-person log", PURPL, PURP, PURP, line=PURP, tsize=9.5, ssize=8)
# links
arrow(s, 6.65, 3.15, 6.65, 2.45, color=GREEN)
arrow(s, 2.6, 3.5, 5.5, 3.5, color=MUTED)
arrow(s, 9.3, 2.0, 7.8, 3.2, color=MUTED)
arrow(s, 10.3, 2.35, 10.3, 2.65, color=MUTED); arrow(s, 11.5, 1.95, 10.6, 2.65, color=MUTED)
arrow(s, 9.3, 4.15, 7.8, 3.6, color=BLUE); arrow(s, 10.3, 3.45, 10.3, 3.75, color=BLUE)
arrow(s, 10.0, 4.85, 10.0, 4.55, color=BLUE)
arrow(s, 7.0, 4.8, 6.9, 4.0, color=PURP); arrow(s, 9.3, 5.2, 7.8, 5.1, color=PURP)
arrow(s, 5.5, 3.7, 2.7, 4.7, color=MUTED); arrow(s, 5.5, 3.8, 4.0, 4.7, color=PURP)
text(s, 0.7, 6.2, 12, 0.4, [("Lines show foreign-key relationships. Colour = domain: ", 10.5, False, MUTED, 0)])
for i,(lab,c) in enumerate([("identity/staff",INK),("compliance",GREEN),("events",GOLD),("participation",BLUE),("engagement",PURP)]):
    box(s, 4.0+i*1.75, 6.22, 0.18, 0.18, c); text(s, 4.22+i*1.75, 6.14, 1.6, 0.3, [(lab, 9.5, False, MUTED, 0)])
page(s)

# 2b — identity
s = slide(); header(s, "Data model · 2 of 8", "Identity & access")
field_table(s, 0.7, 1.95, 6.0, "volunteers", [
    ("id", "uuid · PK"),
    ("auth_user_id", "FK → auth.users (Supabase Auth)"),
    ("email · email_verified", "login id + verification gate"),
    ("gender · date_of_birth", "eligibility + age"),
    ("is_active", "soft enable/disable"),
    ("+ personal / emergency / medical", "address, contact, dietary, medical"),
])
field_table(s, 7.0, 1.95, 5.6, "admins", [
    ("id", "uuid · PK"),
    ("email", "matched to the Auth user"),
    ("role", "super_admin · admin · coordinator"),
    ("is_active · last_login", "access + audit"),
])
callout(s, 7.0, 4.55, 5.6, 1.4, "Auth lives in Supabase",
        "Passwords/sessions are in auth.users. The volunteers / admins tables hold the profile + role; they link by auth_user_id / email. Admin role is also stamped into the JWT (app_metadata).", fill=GOLDL, tc=INK)
page(s)

# 2c — compliance
s = slide(); header(s, "Data model · 3 of 8", "Compliance")
field_table(s, 0.7, 1.95, 6.2, "volunteer_compliance  (1 per volunteer)", [
    ("volunteer_id", "PK · FK → volunteers (unique)"),
    ("lseg_status", "pending · clear · possible_match · high_risk"),
    ("dbs_status", "not_uploaded · pending · verified · rejected"),
    ("overall_status", "pending · approved · rejected (derived)"),
    ("dbs_document_url", "private storage path"),
    ("approved_at / approved_by", "audit"),
])
callout(s, 7.1, 1.95, 5.5, 2.05, "overall_status is trigger-derived",
        "A BEFORE INSERT/UPDATE trigger computes it:\n• LSEG clear → approved\n• LSEG not clear + DBS verified → approved\n• LSEG not clear + DBS rejected → rejected\n• otherwise → pending", fill=GREENL, tc=GREEN)
callout(s, 7.1, 4.2, 5.5, 1.4, "Why derive it",
        "Status is never set by hand, so the app and the gated dashboard always agree. Admin actions just set lseg_status / dbs_status; the trigger does the rest.", fill=GOLDL, tc=INK)
page(s)

# 2d — events & roles
s = slide(); header(s, "Data model · 4 of 8", "Events & roles")
field_table(s, 0.7, 1.95, 6.1, "events", [
    ("id · name · city", "the event"),
    ("status", "draft · published · active · completed · cancelled"),
    ("volunteer_start / event_start", "arrival vs public start"),
    ("event_end / volunteer_end", "public end vs finish (cleanup)"),
    ("created_by", "FK → admins"),
])
field_table(s, 7.0, 1.95, 5.6, "roles  (catalog)", [
    ("id · name", "global, reusable role"),
    ("description", "shown to volunteers on apply"),
    ("is_active", "soft delete"),
], h_each=0.32)
field_table(s, 7.0, 3.75, 5.6, "event_roles  (per event)", [
    ("event_id · role_catalog_id", "FKs → events / roles"),
    ("capacity · gender_restriction", "slots + eligibility"),
], h_each=0.32)
callout(s, 0.7, 5.05, 6.1, 1.05, "Four-time model",
        "Volunteer vs event start/end lets volunteers arrive early / leave late, and drives the points engine.", fill=GOLDL, tc=INK)
page(s)

# 2e — participation
s = slide(); header(s, "Data model · 5 of 8", "Participation")
field_table(s, 0.7, 1.95, 6.1, "event_applications", [
    ("id · volunteer_id · event_id · role_id", "who / where / role"),
    ("status", "applied · confirmed · waitlisted · cancelled · no_show"),
    ("qr_token", "unique check-in QR"),
    ("waitlist_position", "queue order"),
])
field_table(s, 7.0, 1.95, 5.6, "check_ins", [
    ("application_id · volunteer_id · event_id", "FKs"),
    ("station", "entry · exit · setup · cleanup · …"),
    ("scanned_at", "timestamp of scan"),
    ("within_time_window · points_awarded", "punctuality + points"),
])
callout(s, 0.7, 4.65, 11.9, 1.1, "Flow",
        "Apply → confirmed (auto, if capacity) or waitlisted. On the day, an admin scans the QR → a check_in row (entry, then exit). Each scan is graded for punctuality and awards points.", fill=BLUEL, tc=BLUE)
page(s)

# 2f — engagement / points
s = slide(); header(s, "Data model · 6 of 8", "Engagement & points")
field_table(s, 0.7, 1.95, 6.2, "points_transactions  (append-only)", [
    ("volunteer_id · event_id · check_in_id", "source"),
    ("type", "check_in · check_out · manual_bonus · deduction · …"),
    ("amount", "+/- points"),
    ("awarded_by", "FK → admins (null = automatic)"),
])
field_table(s, 7.1, 1.95, 5.5, "points_config  (singleton)", [
    ("check_in / late · check_out / early", "configurable point values"),
    ("grace_minutes", "on-time window"),
], h_each=0.32)
field_table(s, 7.1, 3.75, 5.5, "points_tiers", [
    ("name · min_points", "Bronze 0 · Silver 300 · Gold 800 · Platinum 1500"),
], h_each=0.32)
callout(s, 0.7, 4.6, 6.2, 1.15, "Config-driven engine",
        "Point values, the grace window and tier thresholds live in the DB and are editable in admin Settings — no code change to re-tune.", fill=PURPL, tc=PURP)
page(s)

# 2g — messaging
s = slide(); header(s, "Data model · 7 of 8", "Messaging")
field_table(s, 0.7, 1.95, 6.1, "broadcast_logs / broadcast_recipients", [
    ("channel", "email · sms"),
    ("scope · gender · sender_id", "audience + SMS sender"),
    ("recipient.status", "sent · delivered · opened · failed · invalid"),
    ("resend_message_id / sms_message_id", "provider ids for tracking"),
])
field_table(s, 7.0, 1.95, 5.6, "communications  (per-volunteer log)", [
    ("volunteer_id · channel", "email · sms"),
    ("category", "system · direct"),
    ("subject · status · sent_by", "what was sent, by whom"),
])
callout(s, 7.0, 4.35, 5.6, 1.35, "Unified history",
        "A volunteer's profile shows every message — system, direct and broadcast — by merging communications with broadcast_recipients.", fill=BLUEL, tc=BLUE)
callout(s, 0.7, 4.95, 6.1, 0.95, "email_templates",
        "Every email is an admin-editable template (key, subject, body_html) rendered at send time.", fill=GOLDL, tc=INK)
page(s)

# 2h — system / security tables
s = slide(); header(s, "Data model · 8 of 8", "System & security tables")
field_table(s, 0.7, 1.95, 6.1, "Auth & verification", [
    ("email_verifications", "hashed token · expiry · resend counters"),
    ("admin_otp_challenges", "hashed 6-digit OTP · expiry"),
], h_each=0.34)
field_table(s, 0.7, 3.5, 6.1, "Config & audit", [
    ("sms_config", "default SMS sender id"),
    ("certificate_configs", "per-event certificate template + layout"),
    ("activity_log", "actor · action · metadata (audit trail)"),
], h_each=0.34)
callout(s, 7.0, 1.95, 5.6, 2.0, "Secret-bearing tables are locked down",
        "email_verifications and admin_otp_challenges store only SHA-256 hashes, never the raw token/code. RLS is on with no anon policy → reachable only by the server (service role).", fill=GREENL, tc=GREEN)
callout(s, 7.0, 4.15, 5.6, 1.5, "Single-source config",
        "points_config, points_tiers, sms_config and email_templates make the platform tunable from the admin UI without deploys.", fill=GOLDL, tc=INK)
page(s)

# ════════════════════════════════════════════════════════════════════════
# 3 · AUTH & ONBOARDING
# ════════════════════════════════════════════════════════════════════════
divider("3", "Authentication & Onboarding", "From sign-up to an approved, active volunteer")

# 3a — account model
s = slide(); header(s, "Auth · 1 of 5", "Account model & login paths")
boxlabel(s, 5.15, 1.95, 3.0, 0.85, "auth.users", "Supabase Auth · passwords + JWT", INK, WHITE, RGBColor(0xC9,0xC1,0xB2), tsize=13)
boxlabel(s, 1.2, 3.5, 3.6, 1.0, "volunteers", "email_verified · is_active\nlinked by auth_user_id", BLUEL, BLUE, BLUE, line=BLUE, tsize=13)
boxlabel(s, 8.5, 3.5, 3.6, 1.0, "admins", "role in app_metadata\nlinked by email", GOLDL, INK, MUTED, line=LINE, tsize=13)
arrow(s, 5.6, 2.8, 3.0, 3.5, color=BLUE); arrow(s, 7.7, 2.8, 10.3, 3.5, color=GOLD)
boxlabel(s, 1.2, 4.95, 3.6, 0.8, "Volunteer login", "password + email verified", BLUEL, BLUE, BLUE, line=BLUE, tsize=12, ssize=9.5)
boxlabel(s, 8.5, 4.95, 3.6, 0.8, "Admin login", "password + email OTP (2FA)", GOLDL, INK, MUTED, line=LINE, tsize=12, ssize=9.5)
arrow(s, 3.0, 4.5, 3.0, 4.95, color=BLUE); arrow(s, 10.3, 4.5, 10.3, 4.95, color=GOLD)
callout(s, 1.2, 6.0, 10.9, 0.85, "One identity store, two journeys",
        "Both roles authenticate against Supabase Auth; the profile table determines which experience (and which extra checks) apply.", fill=GOLDL, tc=INK)
page(s)

# 3b — registration + verification
s = slide(); header(s, "Auth · 2 of 5", "Registration → email verification")
flow = [
    ("Register", "3-step form", GOLDL, INK),
    ("Create user\n+ volunteer", "email_verified = false", BLUEL, BLUE),
    ("Issue token", "32-byte random,\nSHA-256 hashed, 1h", GREENL, GREEN),
    ("Email link", "/verify-email?token=…", GOLDL, INK),
    ("Verified", "email_verified = true", GREENL, GREEN),
]
x=0.55; bw=2.3; gap=0.31; y=2.4
for i,(t,sub,fill,tc) in enumerate(flow):
    boxlabel(s, x, y, bw, 1.4, t, sub, fill, tc, tc, line=(tc if fill in (BLUEL,GREENL) else LINE), tsize=12.5, ssize=9.5)
    if i<len(flow)-1: arrow(s, x+bw, y+0.7, x+bw+gap, y+0.7)
    x+=bw+gap
callout(s, 0.55, 4.35, 12.1, 1.0, "Security",
        "Only the hash is stored (a DB leak yields no usable links). One row per volunteer ⇒ a resend invalidates older links. Resend is throttled (60s cooldown + 5/24h) and never reveals whether an address exists.", fill=GREENL, tc=GREEN)
callout(s, 0.55, 5.5, 12.1, 1.05, "Gate",
        "Login is blocked until email_verified = true — the app verifies it explicitly (independent of Supabase's own email-confirm setting). Existing volunteers were grandfathered in.", fill=GOLDL, tc=INK)
page(s)

# 3c — login gate
s = slide(); header(s, "Auth · 3 of 5", "Volunteer login gate")
text(s, 0.7, 2.0, 7, 0.4, [("/api/auth/login checks, in order:", 13, True, INK, 0)])
checks = [
    ("1 · Password", "Supabase signInWithPassword", GREEN),
    ("2 · Volunteer account exists", "linked to this Auth user", GREEN),
    ("3 · Email verified", "else 403 + ‘resend’ option", WARN),
    ("4 · Account active", "is_active = true", GREEN),
]
y=2.55
for t,sub,c in checks:
    box(s, 0.7, y, 0.16, 0.62, c); boxlabel(s, 1.0, y, 6.0, 0.62, t, None, CARD, INK, line=LINE, tsize=12)
    text(s, 4.3, y, 2.6, 0.62, [(sub, 9.5, False, MUTED, 0)], anchor=MSO_ANCHOR.MIDDLE)
    y+=0.74
callout(s, 7.3, 2.55, 5.3, 2.0, "Then — the compliance gate",
        "A verified, active volunteer can sign in, but the dashboard stays gated until overall_status = approved. They see a status screen (under review / upload DBS / rejected) instead.", fill=BLUEL, tc=BLUE)
callout(s, 7.3, 4.75, 5.3, 0.95, "Admins differ",
        "Admin login adds an emailed 6-digit OTP (2FA) before a session is issued.", fill=GOLDL, tc=INK)
page(s)

# 3d — compliance state machine
s = slide(); header(s, "Auth · 4 of 5", "Onboarding state machine", "DBS becomes the fallback when the LSEG check isn't clear")
boxlabel(s, 0.7, 3.1, 2.3, 1.0, "Awaiting review", "just signed up", GOLDL, INK, MUTED, line=LINE, tsize=12.5, ssize=9.5)
boxlabel(s, 3.7, 1.95, 2.6, 0.95, "LSEG clear", "→ approved", GREENL, GREEN, GREEN, line=GREEN, tsize=12.5, ssize=10)
boxlabel(s, 3.7, 3.35, 2.6, 1.0, "Not clear →\nDBS required", "volunteer uploads DBS", WARNL, WARN, WARN, line=WARN, tsize=12, ssize=9)
boxlabel(s, 6.9, 3.35, 2.3, 1.0, "DBS review", "admin reviews", BLUEL, BLUE, BLUE, line=BLUE, tsize=12.5, ssize=9.5)
boxlabel(s, 10.0, 1.95, 2.6, 0.95, "Approved", "dashboard unlocks", GREENL, GREEN, GREEN, line=GREEN, tsize=13, ssize=10)
boxlabel(s, 10.0, 4.55, 2.6, 0.95, "Rejected", "no volunteer status", REDL, RED, RED, line=RED, tsize=13, ssize=10)
arrow(s, 3.0, 3.4, 3.7, 2.4, color=GREEN); arrow(s, 3.0, 3.7, 3.7, 3.8, color=WARN)
arrow(s, 6.3, 3.85, 6.9, 3.85, color=BLUE)
arrow(s, 6.3, 2.4, 10.0, 2.3, color=GREEN)
arrow(s, 9.2, 3.6, 10.0, 2.6, color=GREEN); text(s, 9.0, 2.95, 1.4, 0.3, [("verified", 9, False, GREEN, 0)])
arrow(s, 9.2, 4.1, 10.0, 4.7, color=RED); text(s, 9.0, 4.25, 1.4, 0.3, [("declined", 9, False, RED, 0)])
callout(s, 0.7, 5.95, 11.9, 0.9, "No manual ‘approve’ step",
        "A verified volunteer with a clear LSEG check is approved automatically. The DBS path is the only second route in — and the final decision (approve / reject) is the admin's.", fill=GOLDL, tc=INK)
page(s)

# 3e — gating & RLS recap
s = slide(); header(s, "Auth · 5 of 5", "Session, gating & data protection")
bullets(s, 0.7, 2.0, 6.0, [
    ("Sessions", "are Supabase JWTs in HTTP-only cookies (SameSite=Lax); validated with getUser() on every request."),
    ("Volunteer gating", "is enforced both client-side (the shell shows the status screen) and server-side (event APIs require overall_status = approved)."),
    ("Admin authorization", "checks the JWT role on every /api/admin/* route."),
])
bullets(s, 7.0, 2.0, 5.6, [
    ("RLS is deny-by-default", "— most tables have RLS on with no anon policy, so only the server (service role) can read them."),
    ("Secrets", "(tokens, OTPs) are stored hashed; provider keys are server-only."),
    ("Webhooks", "are signature/token verified before they mutate state."),
])
callout(s, 0.7, 5.4, 11.9, 1.2, "Defense in depth",
        "Identity (Supabase Auth) → app authorization (route guards) → email verification → compliance approval → RLS as the database-level safety net. Each layer is independent, so one mistake doesn't open the door.", fill=GREENL, tc=GREEN)
page(s)

import os
out = os.path.join(os.path.dirname(__file__), "Volunteer-Hub-Infrastructure.pptx")
prs.save(out)
print("Saved:", out, "·", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
